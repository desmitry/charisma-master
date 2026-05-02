import asyncio
import logging
import os
import subprocess
import tempfile
from enum import Enum
from pathlib import Path
from typing import Optional

from celery import shared_task
from celery.signals import worker_process_init
from charisma_schemas import (
    AnalysisResult,
    AnalyzeProvider,
    ConfidenceComponents,
    ConfidenceIndex,
    EvaluationCriteriaReport,
    EvaluationCriterion,
    FillersSummary,
    PersonaRoles,
    TaskStage,
    TaskState,
    TranscribeProvider,
)
from charisma_storage import (
    BUCKET_RESULTS,
    BUCKET_UPLOADS,
    download_file,
    put_object_json,
)

from app.config import settings
from app.logic.competition_research import CompetitionResearchAgent
from app.logic.llm_client import LLMClient
from app.logic.ml_engine import MLEngine
from app.logic.prompts import load_prompts_from_db

logger = logging.getLogger(__name__)


def _run_async(coro):
    loop = asyncio.new_event_loop()
    try:
        asyncio.set_event_loop(loop)
        return loop.run_until_complete(coro)
    finally:
        asyncio.set_event_loop(None)
        loop.close()


def _get_presentation_text(presentation_path: Optional[str]) -> list[str]:
    """Obtaining the presentation text by slides.

    Args:
        presentation_path (Optional[str]): Path to the .pptx or .pdf file.

    Raises:
        RuntimeError: If parsing fails.

    Returns:
        list[str]: Text on each slide.
    """
    if not presentation_path or not Path(presentation_path).exists():
        return []

    file_ext = Path(presentation_path).suffix.lower()

    try:
        if file_ext == ".pptx":
            return _parse_pptx(presentation_path)
        elif file_ext == ".pdf":
            return _parse_pdf(presentation_path)
        else:
            logger.warning(f"Unsupported presentation format: {file_ext}")
            return []
    except Exception as e:
        logger.error(f"Presentation parsing failed: {e}")
        raise


def _parse_pptx(presentation_path: str) -> list[str]:
    """Parse text from PPTX file slide by slide."""
    from pptx import Presentation

    prs = Presentation(presentation_path)
    slides_text = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        slides_text.append("\n".join(slide_text) if slide_text else "")

    return slides_text


def _parse_pdf(presentation_path: str) -> list[str]:
    """Parse text from PDF file page by page."""
    import pdfplumber

    slides_text = []

    with pdfplumber.open(presentation_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            slides_text.append(text.strip() if text else "")

    return slides_text


def _download_to_temp(object_key: str) -> str:
    """Download from SeaweedFS to a temp file, return the path."""
    suffix = Path(object_key).suffix or ".mp4"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp.close()
    download_file(
        settings.seaweedfs_endpoint,
        BUCKET_UPLOADS,
        object_key,
        tmp.name,
        settings.seaweedfs_access_key,
        settings.seaweedfs_secret_key,
    )
    return tmp.name


@worker_process_init.connect
def init_worker(**kwargs):
    logger.info("Preloading prompts from database...")
    try:
        load_prompts_from_db()
    except Exception as e:
        logger.error(f"Failed to load prompts from DB: {e}")

    logger.info("Preloading ML models...")
    try:
        # TODO: Preload all local models
        MLEngine.load_model()
    except Exception as e:
        logger.error(f"Failed to load: {e}")


@shared_task(bind=True)
def process_video_pipeline(  # noqa: C901
    self,
    task_id: str,
    speech_video_key: Optional[str],
    speech_text_key: Optional[str],
    need_text_from_video: bool,
    need_video_analysis: bool,
    presentation_key: Optional[str],
    evaluation_criteria_key: Optional[str],
    evaluation_criteria_preset: Optional[list[dict]],
    persona: str | Enum,
    analyze_provider: str | Enum,
    transcribe_provider: str | Enum,
):
    persona = PersonaRoles(persona)
    analyze_provider = AnalyzeProvider(analyze_provider)
    transcribe_provider = TranscribeProvider(transcribe_provider)

    video_tmp = None
    audio_tmp = None

    full_text = None
    transcript_segments = []
    tempo_data = []
    long_pauses = []

    if need_text_from_video and speech_video_key:
        self.update_state(
            state=TaskState.processing.value,
            meta=TaskStage.transcription.meta,
        )
        if video_tmp is None:
            video_tmp = _download_to_temp(speech_video_key)

        if audio_tmp is None:
            audio_tmp = str(Path(video_tmp).with_suffix(".wav"))

        try:
            MLEngine.extract_audio(
                video_tmp,
                audio_tmp,
            )
            transcript_segments = MLEngine.transcribe(
                audio_tmp,
                provider=transcribe_provider,
            )
            tempo_data = MLEngine.calculate_tempo(
                transcript_segments,
            )
            long_pauses = MLEngine.get_long_pauses(
                transcript_segments,
                threshold=2.0,
            )
            full_text = " ".join([s.text for s in transcript_segments])
        except subprocess.CalledProcessError as e:
            logger.critical("FFmpeg error: %s", e.stderr, exc_info=True)
            raise RuntimeError("Audio extraction failed")
        except Exception as e:
            logger.critical("Global pipeline error: %s", e, exc_info=True)
            raise RuntimeError("Processing failed")

    if speech_text_key and full_text is None:
        text_tmp = _download_to_temp(
            speech_text_key,
        )
        full_text = LLMClient._read_document_text(
            text_tmp,
            Path(text_tmp).suffix,
        )
        os.unlink(
            text_tmp,
        )

    video_metrics = MLEngine.get_empty_video_metrics()
    audio_metrics = MLEngine.get_empty_audio_metrics()
    if need_video_analysis and speech_video_key:
        self.update_state(
            state=TaskState.processing.value,
            meta=TaskStage.video_analisis.meta,
        )
        if video_tmp is None:
            video_tmp = _download_to_temp(speech_video_key)
        if audio_tmp is None:
            audio_tmp = str(Path(video_tmp).with_suffix(".wav"))
            try:
                MLEngine.extract_audio(video_tmp, audio_tmp)
            except subprocess.CalledProcessError as e:
                logger.error(f"FFmpeg audio extraction failed: {e}")
                audio_tmp = None

        try:
            video_metrics = MLEngine.analyze_video(
                video_tmp,
            )
        except Exception as e:
            logger.warning(
                f"Video analysis failed: {e}",
            )

        if audio_tmp:
            try:
                audio_metrics = MLEngine.analyze_audio(
                    audio_tmp,
                )
            except Exception as e:
                logger.warning(
                    f"Audio analysis failed: {e}",
                )

    presentation_tmp = None
    presentation_text = "Текст презентации отсутствует"
    self.update_state(
        state=TaskState.processing.value,
        meta=TaskStage.presentation_text_parsing.meta,
    )
    if presentation_key:
        presentation_tmp = _download_to_temp(presentation_key)
        presentation_text_list = _get_presentation_text(presentation_tmp)
        if presentation_text_list:
            presentation_text = "\n\n".join(
                f"[Слайд: {idx + 1}]\n{text}"
                for idx, text in enumerate(presentation_text_list)
            )

    evaluation_criteria_list = list()
    llm_client = LLMClient()
    self.update_state(
        state=TaskState.processing.value,
        meta=TaskStage.evaluation_criteria_report.meta,
    )
    if evaluation_criteria_preset:
        evaluation_criteria_list = [
            EvaluationCriterion(**c) for c in evaluation_criteria_preset
        ]
    elif evaluation_criteria_key:
        criteria_path = None
        try:
            criteria_path = _download_to_temp(evaluation_criteria_key)
            evaluation_criteria_list = _run_async(
                llm_client.get_evaluation_criteria(criteria_path)
            )
        except Exception as e:
            logger.error(
                "Failed to extract evaluation criteria: %s", e, exc_info=True
            )
            raise RuntimeError("Criteria extraction failed")

        if criteria_path and os.path.exists(criteria_path):
            os.unlink(criteria_path)

    competition_analysis = ""
    if full_text:
        self.update_state(
            state=TaskState.processing.value,
            meta=TaskStage.competition_research.meta,
        )
        try:
            competition_agent = CompetitionResearchAgent(llm_client=llm_client)
            competition_analysis = _run_async(
                competition_agent.run(
                    transcript_text=full_text,
                    presentation_text=presentation_text,
                    provider=analyze_provider,
                )
            )
        except Exception as e:
            logger.warning(
                "Competition research failed, continuing without it: %s",
                e,
                exc_info=True,
            )
            competition_analysis = ""

    llm_speech_report = LLMClient._get_empty_speech_analysis_response()
    self.update_state(
        state=TaskState.processing.value,
        meta=TaskStage.llm_speech_report.meta,
    )

    if full_text:
        try:
            llm_speech_report = _run_async(
                llm_client.analyze_speech(
                    transcript_text=full_text,
                    presentation_text=presentation_text,
                    provider=analyze_provider,
                    persona=persona,
                    competition_analysis=competition_analysis,
                )
            )
            if (
                competition_analysis
                and not llm_speech_report.competition_analysis.strip()
            ):
                llm_speech_report.competition_analysis = competition_analysis
        except Exception as e:
            logger.error("LLM speech analysis failed: %s", e, exc_info=True)
            raise RuntimeError("Speech analysis failed")

    evaluation_criteria_total_score = 0
    evaluation_criteria_max_score = 0
    self.update_state(
        state=TaskState.processing.value,
        meta=TaskStage.llm_criteria_report.meta,
    )
    if full_text and evaluation_criteria_list:
        try:
            evaluation_criteria_list = _run_async(
                llm_client.analyze_with_evalution_criteria(
                    transcript_text=full_text,
                    presentation_text=presentation_text,
                    provider=analyze_provider,
                    evaluation_criteria=evaluation_criteria_list,
                    competition_analysis=competition_analysis,
                )
            )

            evaluation_criteria_total_score = sum(
                c.current_value for c in evaluation_criteria_list
            )
            evaluation_criteria_max_score = sum(
                c.max_value for c in evaluation_criteria_list
            )

            logger.info(
                f"Criteria evaluation completed: "
                f"{evaluation_criteria_total_score}/{evaluation_criteria_max_score}"
            )
        except Exception as e:
            logger.error("Criteria analysis failed: %s", e, exc_info=True)
            raise RuntimeError("Criteria analysis failed")

    # TODO: Document and verify these coefficients.
    total_words = len(full_text.split()) if full_text else 1

    base_filler_count = sum(
        1 for s in transcript_segments for w in s.words if w.is_filler
    )

    filler_ratio = base_filler_count / total_words if total_words > 0 else 0
    filler_score = max(0, 100 - (filler_ratio * 750))

    total_conf = (
        (filler_score * 0.35)
        + (audio_metrics["tone_score"] * 0.25)
        + (audio_metrics["volume_score"] * 0.20)
        + (video_metrics["gaze_score"] * 0.10)
        + (video_metrics["gesture_score"] * 0.10)
    )
    total_conf = min(total_conf, 100)
    total_conf = round(total_conf)
    total_conf_label = MLEngine.get_score_label(total_conf)

    result_data = AnalysisResult(
        task_id=task_id,
        video_path=f"/media/{task_id}.mp4" if speech_video_key else None,
        user_need_video_analysis=need_video_analysis,
        user_need_text_from_video=need_text_from_video,
        transcript=transcript_segments,
        tempo=tempo_data,
        long_pauses=long_pauses,
        fillers_summary=FillersSummary(
            count=base_filler_count, ratio=round(filler_ratio, 4)
        ),
        confidence_index=ConfidenceIndex(
            total=total_conf,
            total_label=total_conf_label,
            components=ConfidenceComponents(
                volume_level=audio_metrics["volume_level"],
                volume_score=round(audio_metrics["volume_score"]),
                volume_label=audio_metrics["volume_label"],
                filler_score=round(filler_score),
                filler_label=MLEngine.get_score_label(filler_score),
                gaze_score=round(video_metrics["gaze_score"]),
                gaze_label=video_metrics["gaze_label"],
                gesture_score=round(video_metrics["gesture_score"]),
                gesture_label=video_metrics["gesture_label"],
                gesture_advice=video_metrics["gesture_advice"],
                tone_score=round(audio_metrics["tone_score"]),
                tone_label=audio_metrics["tone_label"],
            ),
        ),
        speech_report=llm_speech_report,
        evaluation_criteria_report=EvaluationCriteriaReport(
            total_score=evaluation_criteria_total_score,
            max_score=evaluation_criteria_max_score,
            criteria=evaluation_criteria_list,
        ),
        analyze_provider=analyze_provider.value,
        analyze_model=analyze_provider.model_name,
        transcribe_model=transcribe_provider.value,
    )

    put_object_json(
        settings.seaweedfs_endpoint,
        BUCKET_RESULTS,
        f"{task_id}.json",
        result_data.model_dump(),
        settings.seaweedfs_access_key,
        settings.seaweedfs_secret_key,
    )

    return {"status": "completed", "task_id": task_id}
